import streamlit as st
import whisper
import openai
import os
import tempfile
import mammoth
import docx2txt
import pandas as pd
import pptx
from PyPDF2 import PdfReader
from PIL import Image
import pytesseract

# 初期化
if 'summary' not in st.session_state:
    st.session_state.summary = ""
if 'suggestion' not in st.session_state:
    st.session_state.suggestion = ""

openai.api_key = st.secrets["openai_key"]

st.title("GPT Meeting Assistant")

# ユーザー入力
prompt_input = st.text_area("🔧 プロンプト（任意）", placeholder="プロンプトを入力。未入力ならパートナー視点が使われます。")
agenda_files = st.file_uploader("📂 アジェンダファイル（最大10個）", type=["pdf", "doc", "docx", "xls", "xlsx", "ppt", "pptx", "txt", "jpg", "jpeg", "png", "gif"], accept_multiple_files=True)
audio_file = st.file_uploader("🎧 音声ファイルアップロード", type=["mp3", "mp4", "wav", "m4a"])

# 分析開始
if st.button("分析を開始"):
    if audio_file is not None:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as temp_audio:
            temp_audio.write(audio_file.read())
            temp_audio_path = temp_audio.name

        st.write("🔁 Whisperで文字起こし中...")
        transcription = whisper.load_model("medium").transcribe(temp_audio_path, language="ja")
        text = transcription["text"]

        st.write("📑 アジェンダ読み込み中...")
        agenda_texts = []
        for f in agenda_files:
            ext = f.name.split(".")[-1].lower()
            with tempfile.NamedTemporaryFile(delete=False, suffix="." + ext) as tmp:
                tmp.write(f.read())
                path = tmp.name

            try:
                if ext in ["txt"]:
                    with open(path, "r", encoding="utf-8") as file:
                        agenda_texts.append(file.read())
                elif ext in ["pdf"]:
                    pdf = PdfReader(path)
                    content = "\n".join([p.extract_text() or "" for p in pdf.pages])
                    agenda_texts.append(content)
                elif ext in ["docx"]:
                    agenda_texts.append(docx2txt.process(path))
                elif ext in ["doc"]:
                    result = mammoth.convert_to_markdown(open(path, "rb"))
                    agenda_texts.append(result.value)
                elif ext in ["xls", "xlsx"]:
                    df = pd.read_excel(path)
                    agenda_texts.append(df.to_string(index=False))
                elif ext in ["ppt", "pptx"]:
                    pres = pptx.Presentation(path)
                    slides_text = []
                    for slide in pres.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slides_text.append(shape.text)
                    agenda_texts.append("\n".join(slides_text))
                elif ext in ["jpg", "jpeg", "png", "gif"]:
                    img = Image.open(path)
                    text_img = pytesseract.image_to_string(img, lang="jpn")
                    agenda_texts.append(text_img)
            except Exception as e:
                agenda_texts.append(f"(読み込みエラー: {f.name})")

        agenda = "\n\n".join(agenda_texts)
        role_prompt = prompt_input.strip() or (
            "あなたは信頼できるパートナーとして、会話の流れを大切にしながら、"
            "素晴らしい視点には共感を示し、議論に足りない視点には誰もがハッとするような問いをユーモアを交えて提示できます。"
            "問いを出すときは、なぜその問いが必要なのか、答えなければ起こり得るリスクや未来のズレを具体例で示してください。"
            "問いのトーンは柔らかく、それでいて鋭く。「確かに…それ、大事ですね」と思わせる問いを目指してください。"
        )

        # 要約
        st.write("🧠 ChatGPTで要約中...")
        summary_prompt = f"""
以下は会議の文字起こしです。以下の4点を遵守し、事実のみをビジネス向けの丁寧な文章でまとめてください。
1.議論のポイントを漏れなく
2.読みやすく
3.簡潔に
4.構造的に

《アジェンダ》
{agenda}

《文字起こし》
{text}
"""
        summary_response = openai.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": summary_prompt}]
        )
        st.session_state.summary = summary_response.choices[0].message.content

        # 提案
        st.write("💡 ChatGPTで提案中...")
        suggestion_prompt = f"""
以下は会議の文字起こしです。

《アジェンダ》
{agenda}

《文字起こし》
{text}

この内容をもとに、参加者の気づきを促すような
- 改善点
- 問い
- リスクとその回避策

を具体例や理由を添えて、深い共感が得られる形で提案してください。
"""
        suggestion_response = openai.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": role_prompt},
                {"role": "user", "content": suggestion_prompt}
            ]
        )
        st.session_state.suggestion = suggestion_response.choices[0].message.content

        st.success("✅ 分析が完了しました！下からダウンロードできます。")

# 出力表示＆ダウンロード（常時表示）
if st.session_state.summary:
    st.subheader("📄 会議要約")
    st.text_area("会議要約プレビュー", st.session_state.summary, height=200)
    st.download_button(
        label="⬇ 要約をダウンロード",
        data=st.session_state.summary,
        file_name="summary.txt",
        mime="text/plain",
        key="download_summary"
    )

if st.session_state.suggestion:
    st.subheader("💡 ChatGPTの提案")
    st.text_area("提案プレビュー", st.session_state.suggestion, height=200)
    st.download_button(
        label="⬇ 提案をダウンロード",
        data=st.session_state.suggestion,
        file_name="suggestions.txt",
        mime="text/plain",
        key="download_suggestion"
    )
